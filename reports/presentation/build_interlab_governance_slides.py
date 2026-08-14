# -*- coding: utf-8 -*-
"""Append the Interlab-governance section to sae.pptx.

Mirrors build_progress_slides.py's conventions exactly (imported directly,
not re-implemented): manual 16:9 layout, Calibri, warm off-white background,
teal/orange/blue/green/gold accents, figure-first low text density, full
speaker notes via the NOTE_PAYLOADS mechanism.

This section reports methodological and governance progress, not new
scientific results: no ablation has been run. Numbers that come from the
registry are read live at build time (never copied from internship_report.md).
"""

from __future__ import annotations

import shutil

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import build_progress_slides as bp

ROOT = bp.ROOT
DECK = bp.DECK
BACKUP = ROOT / "sae_before_governance_append.pptx"
REPO_ROOT = ROOT.parent.parent
REGISTRY = REPO_ROOT / "registry"

inch = bp.inch
add_bg = bp.add_bg
add_text = bp.add_text
add_multiline = bp.add_multiline
add_title = bp.add_title
add_section_title = bp.add_section_title
add_tag = bp.add_tag
add_card = bp.add_card
add_metric = bp.add_metric
add_arrow = bp.add_arrow
add_slide = bp.add_slide
notes = bp.notes
patch_notes_xml = bp.patch_notes_xml

DARK = bp.DARK
DARK_BG = bp.DARK_BG
GRAY = bp.GRAY
LIGHT_GRAY = bp.LIGHT_GRAY
TEAL = bp.TEAL
TEAL2 = bp.TEAL2
LIGHT_TEAL = bp.LIGHT_TEAL
ORANGE = bp.ORANGE
LIGHT_ORANGE = bp.LIGHT_ORANGE
BLUE = bp.BLUE
LIGHT_BLUE = bp.LIGHT_BLUE
GREEN = bp.GREEN
LIGHT_GREEN = bp.LIGHT_GREEN
GOLD = bp.GOLD
WHITE = bp.WHITE
NEAR_WHITE = RGBColor(0xEF, 0xEF, 0xEF)


def count_registry(artifact_type: str) -> int:
    """Live read of registry/<type>/*.json -- never hardcoded."""
    d = REGISTRY / artifact_type
    if not d.exists():
        return 0
    return len(list(d.glob("*.json")))


def write_notes_markdown(start_slide: int, payloads: list[str]):
    path = ROOT / "sae_governance_speaker_notes.md"
    chunks = ["# Notes de présentation — section gouvernance Interlab\n"]
    for offset, payload in enumerate(payloads):
        slide_num = start_slide + offset
        chunks.append(f"## Diapositive {slide_num}\n")
        chunks.append(payload)
        chunks.append("")
    path.write_text("\n".join(chunks), encoding="utf-8")


def build_deck():
    bp.NOTE_PAYLOADS.clear()
    if not BACKUP.exists():
        shutil.copy2(DECK, BACKUP)

    # Live registry counts (2026-08-03 read; do not hardcode elsewhere).
    n_a1 = count_registry("corpus_manifest")
    n_a3 = count_registry("census_report")
    n_a5 = count_registry("sae_checkpoint")
    n_a6 = count_registry("sae_certificate")
    n_a7 = count_registry("characterization_manifest")
    n_a8 = count_registry("feature_certificate")
    n_a9 = count_registry("intervention_result")
    n_a10 = count_registry("run_card")
    n_a11 = count_registry("claim_report")

    prs = bp.Presentation(str(DECK))
    start_slide = len(prs.slides) + 1

    # Slide 41
    slide = add_slide(prs)
    add_section_title(
        slide,
        "Depuis la dernière rencontre",
        "Le progrès est méthodologique : la chaîne de preuve est corrigée, "
        "la nécessité est spécifiée et pré-enregistrée, l’environnement se verrouille.",
        "Nouvelle section ajoutée au deck existant · gouvernance et méthodologie Interlab",
    )
    add_tag(slide, "CHAÎNE CORRIGÉE", 0.86, 4.15, 2.05, 0.34, fill=LIGHT_TEAL, line=TEAL, color=TEAL, size=10)
    add_tag(slide, "NÉCESSITÉ PRÉ-ENREGISTRÉE", 3.10, 4.15, 2.75, 0.34, fill=LIGHT_ORANGE, line=ORANGE, color=ORANGE, size=10)
    add_tag(slide, "ENVIRONNEMENT EN VERROUILLAGE", 6.05, 4.15, 3.05, 0.34, fill=RGBColor(0xEF, 0xEF, 0xEF), line=GRAY, color=GRAY, size=10)
    notes(
        slide,
        "Ouvrir cette section comme un suivi méthodologique, pas comme un nouveau résultat.",
        "Le progrès rapporté ici rend le prochain résultat scientifique crédible; il n’en constitue pas un lui-même.",
        [
            "Trois avancées : la chaîne de preuve a été corrigée (A8 vient de validate, pas de characterize), "
            "l’expérience de nécessité pour 9056 est entièrement spécifiée et pré-enregistrée, "
            "et l’environnement cluster est en cours de verrouillage pour la reproductibilité (ED-36).",
            "Aucune de ces trois avancées ne produit une nouvelle mesure scientifique : c’est le point de cette slide.",
            "La suite du deck respecte la même règle que la section précédente : établi vs conçu vs non démontré.",
        ],
        "45 s",
        "Slide de rupture sombre, cohérente avec le séparateur ACTE. Les trois étiquettes peuvent être révélées une à une.",
    )

    # Slide 42
    slide = add_slide(prs)
    add_title(
        slide,
        "La chaîne de preuve corrigée",
        "A8 (feature_certificate) est produit par validate, pas par characterize — la correction qui rend la chaîne de certificats valide.",
        "GOUVERNANCE",
    )
    flow = [
        ("A2", "ConceptBattery", "Concepts +\nprobes", TEAL, LIGHT_TEAL),
        ("A3 · SS1", "Census", "Fréquence de\nsurface", TEAL, LIGHT_TEAL),
        ("A7 · SS5", "Characterize", "Index +\ncorpus_max", ORANGE, LIGHT_ORANGE),
        ("A8 · SS6 · G2", "Validate", "Spécificité +\nsensibilité", GREEN, LIGHT_GREEN),
        ("A9 · SS7/8 · G3", "Steer", "Intervention\njugée", BLUE, LIGHT_BLUE),
        ("A9′", "Judge", "Score\nblindé", DARK, NEAR_WHITE),
    ]
    x0 = 0.40
    y = 1.55
    w = 1.28
    h = 1.30
    gap = 0.26
    for idx, (tag, label, desc, color, fill) in enumerate(flow):
        x = x0 + idx * (w + gap)
        add_tag(slide, tag, x + 0.03, y - 0.34, w - 0.06, 0.27, fill=fill, line=color, color=color, size=7.6)
        add_card(slide, x, y, w, h, label, desc, accent=color, fill=WHITE, title_size=11.6, body_size=8.6)
        if idx < len(flow) - 1:
            add_arrow(slide, x + w + 0.03, y + 0.65, x + w + gap - 0.05, y + 0.65, color=GRAY, width=1.5)

    add_text(
        slide,
        "Correction : l’ancien protocole plaçait la certification (G2) sous characterize; "
        "elle appartient à validate. Les deux gates du chemin sont G2 (à A8) et G3 (à l’intervention).",
        0.62, 3.10, 8.75, 0.60, size=12.4, color=DARK, bold=True, align=PP_ALIGN.CENTER,
    )
    add_tag(
        slide,
        f"État réel : A2/A3 vivants (A3={n_a3}); A7={n_a7}, A8={n_a8}, A9={n_a9} — chaîne conçue, encore non peuplée",
        0.85, 3.85, 8.30, 0.42, fill=LIGHT_ORANGE, line=ORANGE, color=ORANGE, size=11,
    )
    notes(
        slide,
        "Corriger publiquement une erreur de protocole avant de s’appuyer dessus.",
        "La chaîne de certification n’est saine que si chaque artefact est produit par le bon stage; cette diapositive documente la correction.",
        [
            "A7 (characterization_manifest) sort de SS5 : il construit l’index et corpus_max, rien de plus.",
            "A8 (feature_certificate) sort de SS6 validate, pas de characterize : spécificité, sensibilité, sélectivité, probe. C’est le GATE G2.",
            "Le GATE G3 est sur l’intervention elle-même (jobs.steer), pas sur le jugement — le jugement (A9′) est une étape distincte, en aval.",
            "Aucun A7/A8/A9 n’existe encore dans le registre : la chaîne est correcte, pas encore exécutée.",
        ],
        "1 min 45 s",
        "Révéler les six blocs de gauche à droite. Terminer sur l’étiquette d’état réel pour ancrer les décomptes vivants du registre.",
    )

    # Slide 43
    slide = add_slide(prs)
    add_title(
        slide,
        "T1.2 — la moitié manquante : la nécessité",
        "La suffisance est démontrée; la nécessité reste à mesurer.",
        "GOUVERNANCE",
    )
    add_card(
        slide, 0.70, 1.35, 4.05, 2.05,
        "Suffisance — établi",
        [
            "Clamp 9056 vers le haut (échelle ≈ 55) sur des prompts neutres.",
            "Induit un contenu fromage mesurable et jugé.",
            "Résultat de la section précédente (steering sweep).",
        ],
        accent=GREEN, fill=LIGHT_GREEN, body_size=12.0,
    )
    add_card(
        slide, 5.05, 1.35, 4.25, 2.05,
        "Nécessité — spécifié, non exécuté",
        [
            "Clamp 9056 vers zéro sur des prompts qui produisent naturellement du fromage.",
            "Hypothèse H1 : le contenu fromage devrait chuter substantiellement.",
            "Aucune génération, aucun jugement produit à ce jour.",
        ],
        accent=ORANGE, fill=LIGHT_ORANGE, body_size=12.0,
    )
    add_tag(slide, "AUCUNE ABLATION EXÉCUTÉE — SUFFISANCE SEULEMENT À CE JOUR", 1.15, 3.62, 7.75, 0.42, fill=RGBColor(0xEF, 0xEF, 0xEF), line=GRAY, color=GRAY, size=11.5)
    add_text(
        slide,
        "Aucun nouveau code : scales_in_max_units: [0.0] via le hook de clamp déjà existant "
        "(_make_clamp_hook, interplab/interventions/hooks.py) — l’ablation s’exprime entièrement en config.",
        0.85, 4.28, 8.30, 0.62, size=12.5, color=DARK, align=PP_ALIGN.CENTER,
    )
    notes(
        slide,
        "Cadrer précisément ce que l’ablation ajoute et ce qu’elle n’ajoute pas encore.",
        "La nécessité est la moitié manquante de la revendication d’identité; sa spécification ne vaut pas son résultat.",
        [
            "Le mécanisme est déjà dans le code : clamper à l’échelle 0.0 revient exactement à mettre la feature à zéro.",
            "C’est l’élément #5 de la feuille de route, encore sans intervention_result associé dans le registre.",
            "Je répète explicitement : aucune ablation n’a encore été exécutée. Cette slide décrit un protocole, pas un résultat.",
        ],
        "1 min 30 s",
        "Révéler les deux cartes côte à côte, puis l’étiquette de garde-fou, puis la phrase mécanisme.",
    )

    # Slide 44
    slide = add_slide(prs)
    add_title(
        slide,
        "Les bras de contrôle",
        "Isoler l’effet de 9056 exige plus qu’un simple avant/après.",
        "GOUVERNANCE",
    )
    rows = [
        ("baseline", "Pas de hook, sur les prompts fromage-éliciteurs.", "Taux naturel de référence.", TEAL, LIGHT_TEAL, False),
        ("steered", "Clamp 9056 → 0.", "L’ablation — signal de nécessité (H1).", GREEN, LIGHT_GREEN, True),
        ("random_feature", "Clamp une feature de contrôle à fréquence appariée → 0.", "Contrôle de spécificité (H2).", ORANGE, LIGHT_ORANGE, False),
        ("prompt_baseline", "Pas de hook, sur des prompts de contrôle de domaine voisin (ED-22).", "Plancher du juge sur du contenu non-fromage.", BLUE, LIGHT_BLUE, False),
    ]
    y = 1.32
    row_h = 0.72
    for arm, action, role, color, fill, bold in rows:
        add_card(slide, 0.62, y, 2.15, row_h, arm, None, accent=color, fill=fill, title_size=13.5)
        add_text(slide, action, 2.92, y + 0.10, 3.55, row_h - 0.15, size=11.2, color=DARK)
        add_text(slide, role, 6.62, y + 0.10, 2.75, row_h - 0.15, size=11.2, color=color, bold=bold)
        y += row_h + 0.10
    add_tag(
        slide,
        "random_direction : dégénéré à l’échelle 0 (≡ baseline) — vérification de cohérence uniquement, ignoré en analyse",
        0.62, y + 0.06, 8.75, 0.36, fill=NEAR_WHITE, line=GRAY, color=GRAY, size=9.8,
    )
    notes(
        slide,
        "Montrer que la comparaison informative repose sur trois bras, pas deux.",
        "Le contrôle de spécificité sépare « cette feature compte » de « annuler n’importe quelle feature dégrade la sortie ».",
        [
            "baseline vs steered donne la nécessité brute; steered vs random_feature donne la spécificité.",
            "random_direction est produit automatiquement en mode claim mais reste dégénéré à l’échelle 0 : à ignorer.",
            "prompt_baseline calibre le plancher du juge, indépendamment de la feature testée.",
        ],
        "1 min 30 s",
        "Révéler les quatre lignes de haut en bas, puis l’étiquette grise en dernier pour ne pas la faire lire comme un cinquième bras informatif.",
    )

    # Slide 45
    slide = add_slide(prs)
    add_title(
        slide,
        "Critères d’acceptation pré-enregistrés",
        "Décidés avant de voir la moindre donnée.",
        "GOUVERNANCE",
    )
    add_card(
        slide, 0.62, 1.28, 4.35, 1.42,
        "H1 · Nécessité",
        "IC bootstrap groupé par prompt (SS9 bootstrap_ci, resampling au niveau prompt) : "
        "baseline − steered entièrement au-dessus de zéro, ET (d de Cohen ≥ 0,5 OU réduction relative ≥ 50 %).",
        accent=GREEN, fill=LIGHT_GREEN, body_size=10.8,
    )
    add_card(
        slide, 5.10, 1.28, 4.30, 1.42,
        "H2 · Spécificité",
        "Effet de spécificité (steered significativement sous random_feature) ET équivalence du contrôle "
        "à ±0,5 avec baseline, sur l’échelle jugée 1–10.",
        accent=ORANGE, fill=LIGHT_ORANGE, body_size=10.8,
    )
    add_card(slide, 0.62, 2.86, 2.82, 0.98, "Verrou 3 graines", "0 / 42 / 123 — aucune moyenne, aucun cherry-picking. H1 et H2 doivent tenir indépendamment sous les trois.", accent=TEAL, fill=LIGHT_TEAL, body_size=9.8)
    add_card(slide, 3.58, 2.86, 2.82, 0.98, "Agrégation", "3 répétitions Lodestar moyennées en un score par prompt avant toute analyse.", accent=BLUE, fill=LIGHT_BLUE, body_size=9.8)
    add_card(slide, 6.54, 2.86, 2.86, 0.98, "Résultat pré-déclaré", "INCONCLUSIVE si le test d’équivalence est sous-puissant — fait partie du protocole, pas une réserve après coup.", accent=GRAY, fill=WHITE, body_size=9.8)
    add_text(
        slide,
        "C’est la diapositive la plus forte scientifiquement de cette section : rien ici n’a été choisi après avoir vu un résultat.",
        0.75, 4.05, 8.10, 0.34, size=12.6, color=DARK, bold=True, align=PP_ALIGN.CENTER,
    )
    notes(
        slide,
        "Faire comprendre pourquoi la pré-inscription protège la revendication de nécessité future.",
        "Un protocole pré-enregistré retire le choix a posteriori du seuil de succès.",
        [
            "bootstrap_ci et effect_size sont les primitives SS9 déjà figées (interplab/stats/stats.py), pas une méthode ad hoc.",
            "H2 est un test en deux parties : il faut à la fois un effet spécifique ET une équivalence du contrôle avec la baseline.",
            "Le verrou à trois graines interdit explicitement le retry sélectif si un seul seed échoue.",
            "INCONCLUSIVE est un résultat valide prévu à l’avance, pas un échec de protocole.",
        ],
        "2 min",
        "Révéler H1 et H2 côte à côte en premier, puis les trois cartes basses, puis la phrase de fermeture.",
    )

    # Slide 46
    slide = add_slide(prs)
    add_title(
        slide,
        "Infrastructure livrée",
        "Ce qui a changé sous le capot depuis la dernière rencontre.",
        "GOUVERNANCE",
    )
    add_card(
        slide, 0.62, 1.35, 2.85, 2.55,
        "Battery v1.1.0",
        [
            "Concept fromage anglais ajouté (12 probes, word_absent vide, status probes_only).",
            "Provenance nommée ED-8 : Mohamed El Yazid — IID.",
            "Golden de tokenisation régénéré (tests/golden/battery_snapshot.json).",
        ],
        accent=TEAL, fill=LIGHT_TEAL, body_size=10.4,
    )
    add_card(
        slide, 3.58, 1.35, 2.85, 2.55,
        "Lanceur de recensement",
        [
            "Publié au commit 9d90ef6 (census SLURM launcher).",
            f"Le recensement A1/A3 tourne maintenant sur Tamia — {n_a1} corpus_manifest, {n_a3} census_report en registre.",
        ],
        accent=BLUE, fill=LIGHT_BLUE, body_size=11.0,
    )
    add_card(
        slide, 6.54, 1.35, 2.85, 2.55,
        "Enveloppe GPU whole-node",
        [
            "Standardisée sur les six lanceurs : census, certify, characterize, steer, train, validate.",
            "h100:4 partout; mem=0 sauf train (100G, par conception) — 4 SAE certifiés (A6) sous ce régime.",
        ],
        accent=ORANGE, fill=LIGHT_ORANGE, body_size=10.6,
    )
    add_text(
        slide,
        "À retenir : les trois éléments sont en production sur main — battery, lanceur de recensement, enveloppe GPU whole-node.",
        0.75, 4.10, 8.15, 0.55, size=12.2, color=DARK, bold=True, align=PP_ALIGN.CENTER,
    )
    notes(
        slide,
        "Présenter l’infrastructure livrée avec la même rigueur que le reste de la section.",
        "Les trois éléments sont vérifiés sur origin/main au moment de cette diapositive, pas supposés.",
        [
            "Battery v1.1.0 : le concept fromage est researcher-authored (ED-8), status probes_only — sensitivity restera "
            "non mesurée tant qu’aucun word_absent n’est fourni, cohérent avec la limite explicite plus loin dans la section.",
            "Le lanceur de recensement est fusionné sur main (commit 9d90ef6).",
            "L’enveloppe GPU whole-node est standardisée sur les six lanceurs, fusionnée — pas seulement préparée sur une branche isolée.",
            f"sae_certificate = {n_a6} en registre, tous produits sous la pile 6.x post-migration.",
        ],
        "1 min 30 s",
        "Révéler les trois cartes de gauche à droite, puis la phrase de synthèse.",
    )

    # Slide 47
    slide = add_slide(prs)
    add_title(
        slide,
        "Reproductibilité — le verrou d’outillage (ED-36)",
        "Un autre laboratoire doit pouvoir reconstruire cet environnement à l’identique.",
        "GOUVERNANCE",
    )
    add_card(slide, 0.62, 1.30, 4.30, 1.35, "Chaîne d’outils entièrement épinglée", "pyproject.toml / uv.lock comme source unique de vérité; export gelé et porteur de hachages (requirements.cluster.txt).", accent=TEAL, fill=LIGHT_TEAL, body_size=10.8)
    add_card(slide, 5.10, 1.30, 4.30, 1.35, "Installation hors ligne", "PIP_NO_INDEX=1, --no-index --require-hashes partout; virtualenv --no-download; aucun contact réseau depuis le nœud de calcul.", accent=BLUE, fill=LIGHT_BLUE, body_size=10.8)
    add_card(slide, 0.62, 2.80, 4.30, 1.35, "Admission par wheel", "Seuls des wheels sont installés directement; un sdist n’entre que comme source d’un wheel dérivé, avec sa propre fiche de provenance vérifiée (build tooling, hachages, cible).", accent=GREEN, fill=LIGHT_GREEN, body_size=10.4)
    add_card(slide, 5.10, 2.80, 4.30, 1.35, "Manifeste d’installation", "Chaque environnement construit enregistre Python, plateforme, CUDA/torch, tous les paquets installés et les hachages du manifeste d’acquisition — vérifiable après coup.", accent=GRAY, fill=WHITE, body_size=10.4)
    notes(
        slide,
        "Présenter le verrouillage d’environnement comme une contribution scientifique, pas seulement opérationnelle.",
        "Sans cette discipline, un futur certificat ne dirait rien de fiable sur quelle bibliothèque a produit ses nombres.",
        [
            "Aucune installation globale n’est permise, sur le cluster comme en local (ED-1, étendu par ED-36).",
            "Le virtualenv est créé --no-download puis pip/setuptools/wheel sont immédiatement remplacés par la version épinglée et vérifiée par hachage du bundle — les paquets embarqués de virtualenv ne servent que d’amorce transitoire.",
            "Un wheel dérivé conserve son sdist source et son hachage correspondant au lock — aucune substitution silencieuse de version n’est possible.",
        ],
        "1 min 45 s",
        "Révéler les quatre cartes en grille 2×2. Prendre le temps sur « admission par wheel », le point le plus technique.",
    )

    # Slide 48
    slide = add_slide(prs)
    add_title(
        slide,
        "LIMITES EXPLICITES",
        "Les garde-fous qui empêchent cette section d’être lue comme une preuve de nécessité.",
        "GOUVERNANCE",
    )
    add_card(
        slide, 0.62, 1.25, 4.30, 1.62,
        "ED-19 · pas de Lodestar en direct",
        "Le conflit numpy-2 non résolu maintient interplab/evaluation/lodestar_adapter.py fermé par conception (fail closed). "
        "Le jugement Stage 2 de l’ablation est bloqué tant que ce n’est pas levé.",
        accent=ORANGE, fill=LIGHT_ORANGE, body_size=10.6,
    )
    add_card(
        slide, 5.10, 1.25, 4.30, 1.62,
        "Stage 1 = préparation seulement",
        "Les trois runs à seed fixe produisent un A9 chacun, mais servent à valider le pipeline. "
        "Explicitement non probants : ils ne comptent pas comme preuve de nécessité et n’alimentent aucun rapport.",
        accent=GRAY, fill=WHITE, body_size=10.6,
    )
    add_card(
        slide, 0.62, 3.02, 4.30, 1.62,
        "Le futur A8 : deux limites avant même d’exister",
        f"characterize/rwu04lpb.yaml utilise judge: \"stub\" (juge de brouillon, pas l’autointerp de production); "
        f"la battery v1 n’a aucun contenu word_absent, donc sensitivity y sera status: \"unavailable\" — "
        f"honnêtement non mesurée, jamais zéro. (feature_certificate en registre : {n_a8}.)",
        accent=BLUE, fill=LIGHT_BLUE, body_size=9.8,
    )
    add_card(
        slide, 5.10, 3.02, 4.30, 1.62,
        "Révision FineWeb jamais capturée",
        "revision: \"unknown\" à l’acquisition. Le corpus est épinglé empiriquement : doc_count 601 369, "
        "token_count 400 000 109, sample_checksum sha256:8312…, avec la révision du tokenizer épinglée exactement (cf98f3b3…).",
        accent=TEAL, fill=LIGHT_TEAL, body_size=9.8,
    )
    notes(
        slide,
        "Fermer la section en rendant explicites toutes les limites qui empêchent un sur-claim.",
        "Chaque limite listée ici est vérifiée dans le code ou le registre au moment de cette diapositive, pas supposée.",
        [
            "ED-19 est la même contrainte numpy qui a mis en pause l’intégration SS8 pendant la migration ED-33 — elle n’a jamais été levée.",
            "Stage 1 vs Stage 2 est une distinction stricte du protocole (docs/ablation_9056_spec.md §6) : seul Stage 2, jugé par Lodestar en direct, compte comme preuve.",
            "Le nom exact du champ de config est judge (pas specificity_judge) — je corrige la formulation pour rester fidèle au schéma A8 réel.",
            "L’absence de revision FineWeb est une limite de provenance historique (ED-8), pas une erreur d’aujourd’hui — le sample_checksum en est le palliatif honnête.",
        ],
        "2 min",
        "Révéler les quatre cartes en grille 2×2, dans l’ordre de lecture. Ne pas accélérer sur la dernière : c’est la garantie de traçabilité du corpus.",
    )

    # A freshly added python-pptx slide has no notesSlide part/relationship
    # until one is accessed; force-create them so patch_notes_xml has a
    # target to overwrite for every new slide.
    for idx in range(start_slide - 1, len(prs.slides)):
        _ = prs.slides[idx].notes_slide

    prs.save(str(DECK))
    write_notes_markdown(start_slide, bp.NOTE_PAYLOADS)
    patch_notes_xml(DECK, start_slide, bp.NOTE_PAYLOADS)
    return start_slide, len(bp.NOTE_PAYLOADS)


if __name__ == "__main__":
    start, count = build_deck()
    print(f"Appended {count} slides starting at slide {start}.")
