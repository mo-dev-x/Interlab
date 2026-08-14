# -*- coding: utf-8 -*-
"""Append the research-progress section to sae.pptx.

The script intentionally mirrors the visual language of the existing deck:
manual 16:9 layouts, Calibri, warm off-white background, dark text, teal and
orange accents, direct figure-first communication, and full speaker notes.
"""

from __future__ import annotations

import shutil
import tempfile
import zipfile
import os
from pathlib import Path
from posixpath import normpath
from xml.etree import ElementTree as ET
from xml.sax.saxutils import escape

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DECK = ROOT / "sae.pptx"
BACKUP = ROOT / "sae_before_progress_append.pptx"

FIG = ROOT / "figures"
DIA = ROOT / "diagrams"

SLIDE_W = 10.0
SLIDE_H = 5.625

BG = RGBColor(0xF6, 0xF5, 0xF4)
DARK = RGBColor(0x1C, 0x1A, 0x18)
DARK_BG = RGBColor(0x2E, 0x2B, 0x28)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xE3, 0xE0, 0xDB)
TEAL = RGBColor(0x0D, 0x73, 0x77)
TEAL2 = RGBColor(0x14, 0x91, 0x9B)
LIGHT_TEAL = RGBColor(0xE6, 0xF4, 0xF4)
ORANGE = RGBColor(0xD6, 0x4A, 0x2A)
LIGHT_ORANGE = RGBColor(0xFA, 0xE9, 0xE4)
BLUE = RGBColor(0x0B, 0x4D, 0xB3)
LIGHT_BLUE = RGBColor(0xE9, 0xF0, 0xFB)
GREEN = RGBColor(0x18, 0x7A, 0x4D)
LIGHT_GREEN = RGBColor(0xE7, 0xF4, 0xEC)
GOLD = RGBColor(0xF0, 0xC2, 0x4B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

NOTE_PAYLOADS: list[str] = []


def inch(value: float):
    return Inches(value)


def add_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_text(run, size=18, color=DARK, bold=False, italic=False):
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color=DARK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    italic: bool = False,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_text(run, size=size, color=color, bold=bold, italic=italic)
    return box


def add_multiline(
    slide,
    lines,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 16,
    color=DARK,
    bold_first: bool = False,
    bullet: bool = False,
    leading: float = 1.0,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.05)
    tf.margin_right = inch(0.05)
    tf.margin_top = inch(0.04)
    tf.margin_bottom = inch(0.04)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3 * leading)
        run = p.add_run()
        run.text = f"• {line}" if bullet else line
        set_text(run, size=size, color=color, bold=(bold_first and idx == 0))
    return box


def add_title(slide, title: str, subtitle: str | None = None, kicker: str | None = None):
    add_text(slide, title, 0.80, 0.32, 8.55, 0.55, size=30, bold=True, color=DARK)
    if subtitle:
        add_text(slide, subtitle, 0.82, 0.90, 8.2, 0.35, size=13.5, color=GRAY)
    if kicker:
        add_tag(slide, kicker, 7.70, 0.32, 1.72, 0.28, fill=LIGHT_TEAL, line=TEAL, color=TEAL, size=9.5)


def add_section_title(slide, title: str, subtitle: str, footer: str | None = None):
    add_bg(slide, DARK_BG)
    add_text(slide, title, 0.80, 1.18, 8.40, 1.08, size=38, color=WHITE, bold=True)
    add_text(slide, subtitle, 0.82, 2.55, 8.20, 0.78, size=18, color=RGBColor(0xD4, 0xD4, 0xD4))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0.82), inch(3.66), inch(1.15), inch(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL2
    line.line.fill.background()
    if footer:
        add_text(slide, footer, 0.82, 4.82, 8.3, 0.30, size=11.5, color=RGBColor(0xB0, 0xB8, 0xC0))


def add_tag(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill=LIGHT_TEAL,
    line=TEAL,
    color=TEAL,
    size=11,
    bold=True,
):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.0)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = inch(0.07)
    tf.margin_right = inch(0.07)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_text(r, size=size, color=color, bold=bold)
    return shp


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str | list[str] | None = None,
    accent=TEAL,
    fill=WHITE,
    title_size=15,
    body_size=12.2,
    number: str | None = None,
):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LIGHT_GRAY
    shp.line.width = Pt(0.9)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(0.06), inch(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    text_x = x + 0.18
    if number:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x + 0.18), inch(y + 0.20), inch(0.34), inch(0.34))
        circ.fill.solid()
        circ.fill.fore_color.rgb = accent
        circ.line.fill.background()
        tf = circ.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = number
        set_text(r, size=11, color=WHITE, bold=True)
        text_x = x + 0.62

    add_text(slide, title, text_x, y + 0.16, w - (text_x - x) - 0.15, 0.33, size=title_size, color=accent, bold=True)
    if body:
        if isinstance(body, str):
            add_text(slide, body, text_x, y + 0.58, w - (text_x - x) - 0.18, h - 0.68, size=body_size, color=DARK)
        else:
            add_multiline(slide, body, text_x, y + 0.55, w - (text_x - x) - 0.18, h - 0.65, size=body_size, color=DARK, bullet=False)
    return shp


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, h: float, color=TEAL, fill=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LIGHT_GRAY
    shp.line.width = Pt(0.8)
    add_text(
        slide,
        value,
        x + 0.08,
        y + 0.08,
        w - 0.16,
        0.25,
        size=18 if len(value) <= 8 else 12.5,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        label,
        x + 0.08,
        y + 0.34,
        w - 0.16,
        h - 0.38,
        size=8.8 if len(label) > 24 else 9.6,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    return shp


def add_image_fit(
    slide,
    image_path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    crop_left: float = 0.0,
    crop_right: float = 0.0,
    crop_top: float = 0.0,
    crop_bottom: float = 0.0,
    border: bool = True,
):
    with Image.open(image_path) as img:
        iw, ih = img.size
    aspect = iw / ih
    box_aspect = w / h
    if aspect >= box_aspect:
        draw_w = w
        draw_h = w / aspect
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * aspect
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    if border:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = LIGHT_GRAY
        bg.line.width = Pt(0.7)
    pic = slide.shapes.add_picture(str(image_path), inch(draw_x), inch(draw_y), width=inch(draw_w), height=inch(draw_h))
    # Avoid python-pptx picture-crop XML here. In this source deck, PowerPoint
    # rejects files containing the generated crop values even though
    # python-pptx can reopen them. Keep the arguments for layout documentation,
    # but insert images by size only.
    return pic


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, width=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def notes(slide, objective: str, message: str, points: list[str], time: str, build: str):
    body = [
        f"Objectif : {objective}",
        f"Message clé : {message}",
        "À dire :",
    ]
    body.extend([f"- {p}" for p in points])
    body.extend([f"Temps estimé : {time}", f"Dévoilement / design : {build}"])
    NOTE_PAYLOADS.append("\n".join(body))


def write_notes_markdown(start_slide: int, payloads: list[str]):
    path = ROOT / "sae_progress_speaker_notes.md"
    chunks = ["# Notes de présentation — section ajoutée\n"]
    for offset, payload in enumerate(payloads):
        slide_num = start_slide + offset
        first = payload.splitlines()[0].replace("Objectif : ", "")
        chunks.append(f"## Diapositive {slide_num}\n")
        chunks.append(payload)
        chunks.append("")
    path.write_text("\n".join(chunks), encoding="utf-8")


def _notes_xml(note_text: str, slide_num: int) -> str:
    paragraphs = []
    for line in note_text.splitlines():
        line = escape(line)
        paragraphs.append(
            '<a:p><a:pPr marL="0" marR="0" lvl="0" indent="0" algn="l" '
            'defTabSz="914400" rtl="0" eaLnBrk="1" fontAlgn="auto" '
            'latinLnBrk="0" hangingPunct="1"><a:lnSpc><a:spcPct val="100000"/>'
            '</a:lnSpc><a:spcBef><a:spcPts val="0"/></a:spcBef>'
            '<a:spcAft><a:spcPts val="0"/></a:spcAft><a:buClrTx/>'
            '<a:buSzTx/><a:buFontTx/><a:buNone/><a:defRPr sz="1200" '
            'dirty="0"/></a:pPr><a:r><a:rPr lang="fr-CA" dirty="0" sz="1200"/>'
            f"<a:t>{line}</a:t></a:r></a:p>"
        )
    if not paragraphs:
        paragraphs.append('<a:p><a:endParaRPr lang="fr-CA"/></a:p>')
    body_xml = "".join(paragraphs)
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr><p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="fr-CA"/></a:p></p:txBody></p:sp><p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/>{body_xml}</p:txBody></p:sp><p:sp><p:nvSpPr><p:cNvPr id="4" name="Slide Number Placeholder 3"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldNum" sz="quarter" idx="10"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:fld id="{{00000000-0000-0000-0000-{slide_num:012d}}}" type="slidenum"><a:rPr lang="fr-CA"/><a:t>{slide_num}</a:t></a:fld><a:endParaRPr lang="fr-CA"/></a:p></p:txBody></p:sp></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:notes>"""


def _notes_target_for_slide(zip_file: zipfile.ZipFile, slide_num: int) -> str:
    rel_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
    root = ET.fromstring(zip_file.read(rel_path))
    ns = {"pr": "http://schemas.openxmlformats.org/package/2006/relationships"}
    for rel in root.findall("pr:Relationship", ns):
        if rel.get("Type", "").endswith("/notesSlide"):
            target = rel.get("Target")
            return normpath(f"ppt/slides/{target}")
    raise RuntimeError(f"No notesSlide relationship found for slide {slide_num}")


def patch_notes_xml(pptx_path: Path, start_slide: int, payloads: list[str]):
    replacements: dict[str, bytes] = {}
    with zipfile.ZipFile(pptx_path, "r") as zin:
        for offset, payload in enumerate(payloads):
            slide_num = start_slide + offset
            target = _notes_target_for_slide(zin, slide_num)
            replacements[target] = _notes_xml(payload, slide_num).encode("utf-8")

        fd, tmp_name = tempfile.mkstemp(suffix=".pptx", dir=str(pptx_path.parent))
        os.close(fd)
        Path(tmp_name).unlink(missing_ok=True)
        with zipfile.ZipFile(tmp_name, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = replacements.get(item.filename)
                if data is None:
                    data = zin.read(item.filename)
                zout.writestr(item, data)

    shutil.move(tmp_name, pptx_path)


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_bg(slide)
    return slide


def build_deck():
    NOTE_PAYLOADS.clear()
    if not BACKUP.exists():
        shutil.copy2(DECK, BACKUP)

    prs = Presentation(str(BACKUP))
    start_slide = len(prs.slides) + 1

    # Slide 25
    slide = add_slide(prs)
    add_section_title(
        slide,
        "Depuis la dernière rencontre",
        "Le projet est passé d’une exploration de steering à une chaîne de preuve instrumentée.",
        "Nouvelle section ajoutée au deck existant · progrès scientifique et infrastructurel",
    )
    add_tag(slide, "INTERLAB", 0.86, 4.15, 1.35, 0.34, fill=LIGHT_TEAL, line=TEAL, color=TEAL, size=10)
    add_tag(slide, "LODESTAR", 2.45, 4.15, 1.55, 0.34, fill=LIGHT_TEAL, line=TEAL, color=TEAL, size=10)
    add_tag(slide, "FEATURE 9056", 4.25, 4.15, 1.72, 0.34, fill=LIGHT_ORANGE, line=ORANGE, color=ORANGE, size=10)
    add_tag(slide, "LIMITES EXPLICITES", 6.25, 4.15, 2.05, 0.34, fill=RGBColor(0xEF, 0xEF, 0xEF), line=GRAY, color=GRAY, size=10)
    notes(
        slide,
        "Ouvrir la nouvelle partie comme un suivi de recherche, sans refaire l’introduction du projet.",
        "La valeur nouvelle est la mise en place d’une chaîne de preuve, pas seulement un nouveau résultat isolé.",
        [
            "Je ne réintroduis pas les SAE ni Golden Gate Claude : on part de ce que la rencontre précédente a déjà établi.",
            "La question implicite est : qu’est-ce qui est maintenant mesurable, certifié, rejetable et limité ?",
            "Les quatre mots-clés annoncent l’équilibre du reste : infrastructure, évaluation, résultat positif, garde-fous.",
        ],
        "45 s",
        "Slide de rupture sombre, cohérente avec les séparateurs existants. Les quatre étiquettes peuvent être révélées de gauche à droite.",
    )

    # Slide 26
    slide = add_slide(prs)
    add_title(slide, "Trois avancées qui changent le statut du projet", "Le progrès principal est méthodologique autant que scientifique.", "ACTE I")
    add_card(
        slide,
        0.78,
        1.48,
        2.72,
        2.55,
        "Infrastructure",
        "Interlab donne une identité vérifiable aux points de contrôle, certificats et artefacts.",
        accent=TEAL,
        fill=LIGHT_TEAL,
        number="1",
    )
    add_card(
        slide,
        3.66,
        1.48,
        2.72,
        2.55,
        "Évaluation",
        "Lodestar transforme les balayages en frontières cohérence-pertinence et points opératoires auditables.",
        accent=BLUE,
        fill=LIGHT_BLUE,
        number="2",
    )
    add_card(
        slide,
        6.54,
        1.48,
        2.72,
        2.55,
        "Résultat",
        "La feature 9056 reproduit quantitativement un effet d’identité sous steering sur Qwen2.5-14B-Instruct.",
        accent=ORANGE,
        fill=LIGHT_ORANGE,
        number="3",
    )
    add_text(slide, "À retenir : le projet produit maintenant des preuves, pas seulement des générations intéressantes.", 1.05, 4.58, 7.9, 0.36, size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    notes(
        slide,
        "Donner le résumé exécutif avant d’entrer dans les preuves.",
        "Les avancées se renforcent mutuellement : l’ingénierie rend la science plus crédible.",
        [
            "Interlab et Lodestar ne sont pas des annexes : ils définissent ce qui peut être revendiqué proprement.",
            "Le résultat 9056 est important parce qu’il est évalué dans ce cadre, pas parce qu’il produit une phrase amusante.",
            "Je vais montrer d’abord l’infrastructure qui rend les chiffres interprétables, puis les résultats.",
        ],
        "1 min 15 s",
        "Révéler les trois cartes dans l’ordre infrastructure → évaluation → résultat. La phrase finale sert de transition.",
    )

    # Slide 27
    slide = add_slide(prs)
    add_title(slide, "La chaîne expérimentale est maintenant instrumentée", "Chaque étape critique a un rôle explicite dans la production de preuve.", "ACTE II")
    add_image_fit(slide, DIA / "Figure1_v3.png", 0.45, 1.18, 9.10, 3.30, crop_bottom=0.05)
    add_metric(slide, "9", "étapes du corpus au rapport", 0.80, 4.67, 1.55, 0.60, color=TEAL)
    add_metric(slide, "A1/A3/A5/A6/A10", "artefacts vivants dans Interlab", 2.55, 4.67, 2.50, 0.60, color=BLUE)
    add_metric(slide, "SS8 externe", "Lodestar exercé en mode autonome", 5.25, 4.67, 2.10, 0.60, color=ORANGE)
    add_metric(slide, "moins d’ambiguïté", "script, artefact, statut", 7.58, 4.67, 1.68, 0.60, color=TEAL)
    notes(
        slide,
        "Montrer le pipeline comme structure de preuve plutôt que comme liste de scripts.",
        "La nouveauté est que les résultats sont désormais attachés à des étapes, statuts et artefacts identifiables.",
        [
            "La figure n’est pas là pour expliquer chaque étape : elle sert à montrer que le pipeline n’est plus implicite.",
            "Les étapes 1 à 3 alimentent la certification; les étapes 6 et 7 produisent les résultats jugés.",
            "Point important : Lodestar a été exercé de façon autonome, alors que son intégration Interlab complète reste un statut distinct.",
        ],
        "1 min 30 s",
        "Commencer par la ligne supérieure, puis pointer les trois zones : certification, steering, évaluation jugée. Les métriques du bas peuvent apparaître ensuite.",
    )

    # Slide 28
    slide = add_slide(prs)
    add_title(slide, "Certifier les SAE avant de chercher des features", "La certification devient un seuil de santé, pas un verdict de monosemanticité.", "ACTE II")
    add_image_fit(slide, FIG / "fig_sae_certification.png", 0.58, 1.28, 8.85, 2.38, border=True)
    add_metric(slide, "4", "SAE certifiés", 0.82, 4.00, 1.48, 0.72, color=TEAL)
    add_metric(slide, "≥ 0,9785", "CE récupéré", 2.52, 4.00, 1.70, 0.72, color=BLUE)
    add_metric(slide, "≤ 0,0020", "fraction de features mortes", 4.45, 4.00, 1.95, 0.72, color=GREEN)
    add_metric(slide, "rwu04lpb", "bande amber, L28×32, résultat 9056", 6.62, 4.00, 2.20, 0.72, color=ORANGE)
    add_text(slide, "La certification mesure l’instrument; la qualité locale d’une feature se démontre par triangulation.", 0.94, 4.93, 8.1, 0.30, size=12.5, color=GRAY, align=PP_ALIGN.CENTER)
    notes(
        slide,
        "Établir que les résultats scientifiques reposent sur des SAE passés par un contrôle de santé.",
        "Le checkpoint utilisé pour 9056 est sain mais pas le meilleur selon la bande, donc la certification ne remplace pas la validation de feature.",
        [
            "Le SAE rwu04lpb est amber, avec CE récupéré 0,9884 et dead-feature fraction 0,0008.",
            "Le seul green est o1cx1dow, mais il n’est pas celui qui porte le résultat principal.",
            "C’est précisément pourquoi il faut distinguer santé globale du SAE et qualité locale d’une feature.",
        ],
        "1 min 40 s",
        "Révéler la figure d’abord, puis les quatre métriques. Insister sur le dernier encadré rwu04lpb.",
    )

    # Slide 29
    slide = add_slide(prs)
    add_title(slide, "Interlab : chaîne de preuve", "De la question scientifique à une revendication vérifiable.", "ACTE II")
    add_card(
        slide,
        0.72,
        1.28,
        8.55,
        0.92,
        "Question scientifique",
        "La couche 28 contient-elle une feature « fromage » proprement steerable ?",
        accent=TEAL,
        fill=LIGHT_TEAL,
        title_size=13.5,
        body_size=12.0,
    )

    flow = [
        ("A1/A3", "Corpus\n+ concepts", "Quelle distribution\net quels probes ?", TEAL, LIGHT_TEAL),
        ("A5/A6 · G1", "Checkpoint\ncertifié", "L’instrument\nest-il sain ?", GREEN, LIGHT_GREEN),
        ("A7/A8 · G2", "Feature\ncandidate", "Est-elle sélective\net spécifique ?", ORANGE, LIGHT_ORANGE),
        ("A9 · SS8", "Intervention\njugée", "Le steering est-il\npropre et évalué ?", BLUE, LIGHT_BLUE),
        ("A11 · G4", "Claim\nreport", "La revendication\nest-elle certifiable ?", DARK, RGBColor(0xEF, 0xEF, 0xEF)),
    ]
    x0 = 0.64
    y = 2.56
    w = 1.54
    h = 1.23
    gap = 0.42
    for idx, (tag, label, question, color, fill) in enumerate(flow):
        x = x0 + idx * (w + gap)
        add_tag(slide, tag, x + 0.12, y - 0.33, w - 0.24, 0.26, fill=fill, line=color, color=color, size=8.4)
        add_card(
            slide,
            x,
            y,
            w,
            h,
            label,
            question,
            accent=color,
            fill=WHITE,
            title_size=12.4,
            body_size=9.3,
        )
        if idx < len(flow) - 1:
            add_arrow(slide, x + w + 0.04, y + 0.62, x + w + gap - 0.08, y + 0.62, color=GRAY, width=1.6)

    add_text(
        slide,
        "Idée clé : Interlab ne certifie pas une intuition; il vérifie une chaîne d’artefacts.",
        0.80,
        4.10,
        8.40,
        0.30,
        size=13.2,
        color=DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_tag(
        slide,
        "État réel ici : G1 vivant + Lodestar autonome; A8/A9/A11 restent conçus, non peuplés.",
        1.22,
        4.65,
        7.55,
        0.38,
        fill=LIGHT_ORANGE,
        line=ORANGE,
        color=ORANGE,
        size=10.4,
    )
    notes(
        slide,
        "Donner une intuition simple d’Interlab avant la vue d’architecture détaillée.",
        "Interlab est une chaîne cible de certification : chaque passage de la question vers la revendication doit laisser un artefact vérifiable.",
        [
            "Je présente cette slide comme l’architecture cible, pas comme un état déjà entièrement peuplé.",
            "La question de départ devient progressivement des artefacts : corpus, checkpoint, certificat SAE, feature, intervention jugée, puis claim report.",
            "Dans le projet actuel, la partie vivante est G1 avec A1/A3/A5/A6/A10; les jugements existent via Lodestar autonome, mais ne sont pas encore repliés dans A9/A11.",
        ],
        "2 min 15 s",
        "Révéler le bandeau question, puis les cinq blocs de gauche à droite. Terminer par l’étiquette d’état réel pour éviter tout sur-claim.",
    )

    # Slide 29
    slide = add_slide(prs)
    add_title(slide, "Interlab : frontière actuelle au jalon G1", "L’état du laboratoire est utile parce qu’il est explicitement borné.", "ACTE II")
    add_card(
        slide,
        0.70,
        1.35,
        2.55,
        1.05,
        "Vivant",
        "A1, A3, A5, A6 et A10 sont peuplés par des artefacts réels.",
        accent=GREEN,
        fill=LIGHT_GREEN,
    )
    add_card(
        slide,
        0.70,
        2.55,
        2.55,
        1.05,
        "Conçu",
        "A8, A9, A11 et l’intégration SS8 ont des schémas complets mais non peuplés.",
        accent=ORANGE,
        fill=LIGHT_ORANGE,
    )
    add_card(
        slide,
        0.70,
        3.75,
        2.55,
        1.05,
        "À ne pas confondre",
        "Les jugements existent bien : ils viennent de Lodestar en mode autonome pour cette campagne.",
        accent=BLUE,
        fill=LIGHT_BLUE,
    )
    add_image_fit(slide, DIA / "Figure4_v2.png", 3.45, 1.18, 6.05, 3.90, crop_bottom=0.08)
    notes(
        slide,
        "Présenter la maturité d’Interlab sans sur-vendre la partie encore non exercée.",
        "Le jalon G1 est vivant; la chaîne complète de rapport de revendication est conçue mais pas encore validée en conditions réelles.",
        [
            "La valeur produite est déjà importante : les checkpoints et certificats sont content-addressed et citables par hash.",
            "Mais il faut garder l’honnêteté du statut : feature_certificate, intervention_result et claim_report ne sont pas encore peuplés.",
            "Cette distinction protège les résultats Section 3 : ils sont jugés par Lodestar, mais pas encore repliés dans A9/A11 Interlab.",
        ],
        "1 min 45 s",
        "Commencer par les trois cartes à gauche, puis utiliser la figure pour situer « You are here » au jalon G1.",
    )

    # Slide 30
    slide = add_slide(prs)
    add_title(slide, "Lodestar ferme la boucle d’évaluation", "Le choix de l’échelle devient une décision mesurée, pas une inspection à l’œil.", "ACTE II")
    add_image_fit(slide, DIA / "Figure5_v2.png", 0.55, 1.05, 4.85, 4.20, crop_bottom=0.07)
    add_card(slide, 5.75, 1.22, 3.55, 0.78, "Frontière", "cohérence × pertinence au lieu d’un score isolé", accent=BLUE, fill=LIGHT_BLUE)
    add_card(slide, 5.75, 2.13, 3.55, 0.78, "Point opératoire", "maximiser la pertinence sous cohérence ≥ 5", accent=TEAL, fill=LIGHT_TEAL)
    add_card(slide, 5.75, 3.04, 3.55, 0.78, "Cache + budget", "coûts bornés, re-jugement évité", accent=GREEN, fill=LIGHT_GREEN)
    add_card(slide, 5.75, 3.95, 3.55, 0.78, "Revue humaine", "le rapport HTML a révélé un défaut de frontière", accent=ORANGE, fill=LIGHT_ORANGE)
    notes(
        slide,
        "Mettre Lodestar au même niveau que les résultats, comme contribution d’évaluation.",
        "Lodestar transforme le steering en boucle expérimentale : générer, juger, inspecter, raffiner.",
        [
            "Le cas fromage illustre la boucle : le balayage 40-150 a montré un trou entre les échelles 40 et 60; un balayage ciblé 45/50/55 a localisé l’optimum.",
            "Le cas Montréal illustre l’autre moitié : une estimation manuelle à 90 a été corrigée par la frontière jugée vers 80.",
            "Le rapport HTML n’est pas une sortie décorative; il a servi à détecter le bug de sweep_hash.",
        ],
        "1 min 45 s",
        "Figure à gauche comme vue système. Révéler les cartes de droite une par une pour éviter de lire tout le schéma.",
    )

    # Slide 31
    slide = add_slide(prs)
    add_title(slide, "Les rapports HTML deviennent l’interface d’analyse", "Un résultat peut être inspecté depuis la frontière jusqu’à chaque génération jugée.", "ACTE II")
    add_image_fit(slide, DIA / "Figure6.png", 0.45, 1.16, 9.10, 3.78, crop_bottom=0.08)
    add_tag(slide, "frontière", 2.96, 1.05, 1.05, 0.28, fill=RGBColor(0x19, 0x18, 0x15), line=GOLD, color=GOLD, size=9)
    add_tag(slide, "générations", 6.70, 1.05, 1.30, 0.28, fill=RGBColor(0x19, 0x18, 0x15), line=GOLD, color=GOLD, size=9)
    add_tag(slide, "validation", 0.76, 3.60, 1.15, 0.28, fill=RGBColor(0x19, 0x18, 0x15), line=GOLD, color=GOLD, size=9)
    add_text(slide, "Les chiffres du rapport se relient à des jugements, pas à des impressions.", 1.05, 5.08, 7.90, 0.28, size=13.2, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    notes(
        slide,
        "Montrer le format concret de l’évidence produite.",
        "Les résultats sont consultables, auditables et re-analysables dans un fichier autonome.",
        [
            "Le rapport contient la frontière, les points opératoires, le contrôle, la validation du juge et les générations individuelles.",
            "La correction du bug sweep_hash a pu être régénérée depuis le cache sans coût de jugement supplémentaire.",
            "C’est une contribution d’ingénierie directement scientifique : elle change la manière de décider si un sweep est crédible.",
        ],
        "1 min 30 s",
        "Présenter le composite comme une interface. Les trois étiquettes peuvent apparaître successivement : frontière, générations, validation.",
    )

    # Slide 32
    slide = add_slide(prs)
    add_title(slide, "Résultat principal : feature 9056", "Reproduction quantitative d’un effet d’identité sous intervention SAE sur Qwen2.5-14B-Instruct.", "ACTE III")
    add_metric(slide, "55", "échelle sélectionnée", 0.78, 1.24, 1.05, 0.66, color=ORANGE, fill=LIGHT_ORANGE)
    add_metric(slide, "5,38", "cohérence", 1.98, 1.24, 1.05, 0.66, color=BLUE, fill=LIGHT_BLUE)
    add_metric(slide, "5,50", "pertinence concept", 3.18, 1.24, 1.24, 0.66, color=ORANGE, fill=LIGHT_ORANGE)
    add_metric(slide, "3,13", "adhérence consigne", 4.57, 1.24, 1.18, 0.66, color=TEAL, fill=LIGHT_TEAL)
    add_tag(slide, "suffisance seulement", 6.05, 1.40, 1.85, 0.34, fill=RGBColor(0xEF, 0xEF, 0xEF), line=GRAY, color=GRAY, size=10)
    add_image_fit(slide, FIG / "fig2_cheese_sweep_judged.png", 0.75, 2.06, 8.50, 2.95, border=True)
    notes(
        slide,
        "Présenter le résultat scientifique central sans l’exagérer.",
        "La feature 9056 suffit à induire un effet d’identité conceptuelle tout en restant au-dessus du plancher de cohérence.",
        [
            "Le point opératoire est l’échelle 55, choisie par Lodestar avec la contrainte cohérence ≥ 5.",
            "La pertinence maximale arrive à des échelles plus élevées, mais la cohérence chute ou devient moins contrôlée.",
            "Ne pas dire que la feature est nécessaire : aucune ablation n’a encore montré que retirer 9056 élimine l’effet.",
        ],
        "2 min",
        "Révéler d’abord les métriques, puis la courbe. L’étiquette « suffisance seulement » doit rester visible comme garde-fou.",
    )

    # Slide 33
    slide = add_slide(prs)
    add_title(slide, "Pourquoi 55, et pas simplement « plus haut » ?", "Le point retenu maximise la pertinence sous une contrainte de lisibilité.", "ACTE III")
    add_image_fit(slide, FIG / "fig3_cheese_mid_judged.png", 0.58, 1.22, 5.78, 3.70, border=True)
    add_card(slide, 6.65, 1.34, 2.72, 0.86, "Règle", "max concept relevance avec cohérence ≥ 5", accent=BLUE, fill=LIGHT_BLUE)
    add_card(slide, 6.65, 2.34, 2.72, 0.86, "Scale 40", "cohérent, mais effet conceptuel faible", accent=GRAY, fill=WHITE)
    add_card(slide, 6.65, 3.34, 2.72, 0.86, "Scale 60", "pertinence plus forte, mais sous le plancher", accent=ORANGE, fill=LIGHT_ORANGE)
    add_text(slide, "La décision expérimentale est la frontière, pas le meilleur exemple isolé.", 0.92, 5.08, 8.15, 0.28, size=13.2, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    notes(
        slide,
        "Faire comprendre le choix de l’échelle comme une décision méthodologique.",
        "L’échelle 55 n’est pas un maximum arbitraire : c’est le meilleur compromis sous une contrainte définie avant l’interprétation.",
        [
            "Cette slide justifie pourquoi le résultat n’est pas choisi à l’œil.",
            "Le balayage ciblé répond à une question précise soulevée par la première frontière : le bon point est entre 50 et 60.",
            "C’est aussi une démonstration de la valeur de Lodestar : il guide le prochain sweep.",
        ],
        "1 min 30 s",
        "Commencer par la courbe. Ajouter ensuite les trois cartes règle, scale 40, scale 60.",
    )

    # Slide 34
    slide = add_slide(prs)
    add_title(slide, "La candidate fiable est sélectionnée par convergence", "Trois mesures indépendantes donnent le même classement.", "ACTE III")
    add_text(slide, "9056  >  47735  >  44189", 0.82, 1.06, 4.00, 0.42, size=21, color=TEAL, bold=True)
    add_text(slide, "fromage  >  UNESCO  >  Eurovision rejetée", 4.95, 1.12, 4.20, 0.30, size=12.8, color=GRAY, align=PP_ALIGN.RIGHT)
    add_image_fit(slide, DIA / "Figure2_v3.png", 0.53, 1.48, 8.95, 3.63, crop_bottom=0.08)
    notes(
        slide,
        "Présenter la triangulation comme contribution scientifique centrale.",
        "La méthodologie ne se contente pas de trouver une feature positive; elle rejette aussi les candidates faibles.",
        [
            "Les trois colonnes sont indépendantes : étiquette de caractérisation, steering jugé, contrôle à taux apparié.",
            "9056 gagne par intégration : pas seulement pertinence, mais adhérence au prompt et naturalité.",
            "44189 est importante comme négatif propre : elle est rejetée par les trois signaux, ce qui rend la méthode falsifiable.",
        ],
        "2 min",
        "Dévoilement recommandé par colonnes : 1 caractérisation, 2 jugement, 3 contrôle, puis le classement final.",
    )

    # Slide 35
    slide = add_slide(prs)
    add_title(slide, "La sélectivité confirme le classement", "Les contrôles à taux apparié évitent de confondre rareté et monosemanticité.", "ACTE III")
    add_image_fit(slide, FIG / "fig_feature_selectivity.png", 0.58, 1.20, 6.35, 2.72, border=True)
    add_card(slide, 7.18, 1.20, 2.13, 0.78, "9056", "max 47,5 · 14,5× médiane · n=1003", accent=GREEN, fill=LIGHT_GREEN)
    add_card(slide, 7.18, 2.12, 2.13, 0.78, "47735", "propre, mais moins bien intégré au prompt", accent=TEAL, fill=LIGHT_TEAL)
    add_card(slide, 7.18, 3.04, 2.13, 0.78, "44189", "max 8,5, sous son contrôle à taux apparié", accent=ORANGE, fill=LIGHT_ORANGE)
    add_text(slide, "Garde-fou : characterize_lite est une preuve de rapport, pas encore un certificat Interlab A7/A8.", 0.90, 4.54, 8.15, 0.34, size=12.5, color=GRAY, align=PP_ALIGN.CENTER)
    notes(
        slide,
        "Montrer que la sélection de 9056 est soutenue par les activations corpus, pas seulement par le steering.",
        "9056 se distingue par la force d’activation et le taux de firing relatif; Eurovision échoue même contre un contrôle de taux similaire.",
        [
            "Le contrôle à taux apparié répond à la question : est-ce seulement une feature qui s’active rarement ou souvent ?",
            "Pour 9056, le max est plus de deux fois celui du contrôle associé.",
            "Pour 44189, le contrôle dépasse la candidate, ce qui rend son rejet méthodologiquement solide.",
            "Nuance importante : ces nombres viennent de characterize_lite, pas encore de la lane Interlab complète.",
        ],
        "1 min 30 s",
        "Révéler d’abord les deux bar charts, puis les trois cartes. La phrase de garde-fou doit rester au bas de la slide.",
    )

    # Slide 36
    slide = add_slide(prs)
    add_title(slide, "Les contrastes rendent la méthode falsifiable", "UNESCO et Eurovision montrent que « pertinence » seule ne suffit pas.", "ACTE III")
    add_image_fit(slide, FIG / "fig8_unesco_judged.png", 0.55, 1.18, 4.32, 2.82, border=True)
    add_image_fit(slide, FIG / "fig9_eurovision_judged.png", 5.12, 1.18, 4.32, 2.82, border=True)
    add_card(slide, 0.70, 4.22, 4.05, 0.72, "47735 · UNESCO", "pertinence haute, mais adhérence et intégration faibles : la candidate prend le dessus", accent=TEAL, fill=LIGHT_TEAL)
    add_card(slide, 5.27, 4.22, 4.05, 0.72, "44189 · Eurovision", "cohérence au plancher et rejet par les trois signaux", accent=ORANGE, fill=LIGHT_ORANGE)
    notes(
        slide,
        "Utiliser les deux candidates non principales comme validation de la méthode de sélection.",
        "Une bonne feature de steering ne se juge pas seulement à la pertinence; elle doit préserver le prompt et l’intégration.",
        [
            "UNESCO est instructive : elle peut produire une pertinence conceptuelle forte, mais au prix d’une réponse moins intégrée.",
            "Eurovision est encore plus clair : elle est faible selon la caractérisation, selon le contrôle et selon Lodestar.",
            "Cette slide évite un récit trop positif : la méthode fonctionne aussi parce qu’elle sait dire non.",
        ],
        "1 min 45 s",
        "Afficher les deux courbes ensemble, puis révéler les deux cartes de synthèse.",
    )

    # Slide 37
    slide = add_slide(prs)
    add_title(slide, "Multilingue : attention à l’unité d’analyse", "Le chevauchement top-20 ne prouve pas l’existence d’une feature unique.", "ACTE III")
    add_image_fit(slide, FIG / "fig_multilingual_overlap.png", 0.48, 1.22, 9.05, 2.70, border=True)
    add_metric(slide, "13/20", "world_cup partagé", 0.92, 4.28, 1.45, 0.66, color=GREEN, fill=LIGHT_GREEN)
    add_metric(slide, "12/20", "quebec partagé", 2.62, 4.28, 1.45, 0.66, color=TEAL, fill=LIGHT_TEAL)
    add_metric(slide, "10/20", "poutine partagé", 4.32, 4.28, 1.45, 0.66, color=ORANGE, fill=LIGHT_ORANGE)
    add_metric(slide, "4/20", "couscous partagé", 6.02, 4.28, 1.45, 0.66, color=GRAY, fill=WHITE)
    add_tag(slide, "ensemble de features ≠ feature monosemantique", 7.72, 4.44, 1.82, 0.35, fill=RGBColor(0xEF, 0xEF, 0xEF), line=GRAY, color=GRAY, size=8.5)
    notes(
        slide,
        "Présenter le résultat multilingue sans créer de contradiction artificielle avec le négatif poutine.",
        "Poutine montre un chevauchement d’ensembles de features, mais cela ne veut pas dire qu’une feature poutine propre existe.",
        [
            "L’unité de mesure est le top-20 features par concept et langue.",
            "Le Jaccard moyen de poutine est 0,51, ce qui indique un chevauchement de voisinage conceptuel.",
            "Ce résultat est compatible avec l’échec à isoler une feature unique après 16+ tentatives.",
            "L’ordre world_cup > quebec > poutine > couscous reste interprétatif; aucun recensement de prévalence n’a été validé.",
        ],
        "1 min 30 s",
        "Révéler les métriques du bas après la heatmap. Finir sur l’étiquette grise pour fixer le garde-fou.",
    )

    # Slide 38
    slide = add_slide(prs)
    add_title(slide, "Les résultats négatifs ont amélioré la méthode", "Ils identifient des mécanismes, pas seulement des échecs expérimentaux.", "ACTE III")
    add_card(slide, 0.70, 1.26, 2.82, 0.78, "Poutine", "16+ tentatives, pas de feature propre; couverture corpus probablement limitante", accent=ORANGE, fill=LIGHT_ORANGE)
    add_card(slide, 0.70, 2.18, 2.82, 0.78, "Montréal / Québec", "feature initialement positive, ensuite corrigée comme enchevêtrement bilingue", accent=BLUE, fill=LIGHT_BLUE)
    add_card(slide, 0.70, 3.10, 2.82, 0.78, "Base → instruct", "un SAE base ne transfère pas automatiquement à la géométrie instruct", accent=TEAL, fill=LIGHT_TEAL)
    add_card(slide, 0.70, 4.02, 2.82, 0.78, "Haute échelle", "sur feature enchevêtrée, la fluence casse avant un thème stable", accent=GRAY, fill=WHITE)
    add_image_fit(slide, FIG / "fig11_montreal_judged.png", 3.86, 1.35, 5.55, 3.45, border=True)
    notes(
        slide,
        "Faire des négatifs une source de crédibilité et de décisions futures.",
        "Les négatifs expliquent où la méthode doit être bornée : corpus, enchevêtrement, géométrie instruct, échelle élevée.",
        [
            "Poutine n’est pas un simple manque de chance : la couverture du corpus et la composition conceptuelle sont le mécanisme plausible.",
            "Montréal/Québec est important parce que l’équipe a corrigé son propre résultat positif.",
            "Le non-transfert base→instruct est un résultat méthodologique, même s’il repose sur un cas unique.",
            "La courbe Montréal rappelle que les features enchevêtrées ne reproduisent pas le régime obsédé mais fluide de Golden Gate Claude.",
        ],
        "2 min",
        "Révéler les quatre cartes comme une liste de leçons, puis montrer la courbe Montréal comme preuve visuelle du compromis défavorable.",
    )

    # Slide 39
    slide = add_slide(prs)
    add_title(slide, "Ce qui est établi aujourd’hui", "La conclusion actuelle est une carte de maturité, pas une fermeture du projet.", "ACTE IV")
    add_card(
        slide,
        0.72,
        1.40,
        2.72,
        3.25,
        "Établi",
        [
            "9056 suffit à induire l’effet",
            "triangulation : 9056 > 47735 > 44189",
            "Lodestar exercé sur les points opératoires",
            "Interlab vivant jusqu’au jalon G1",
        ],
        accent=GREEN,
        fill=LIGHT_GREEN,
        body_size=11.4,
    )
    add_card(
        slide,
        3.64,
        1.40,
        2.72,
        3.25,
        "À interpréter prudemment",
        [
            "globalité multilingue qualitative",
            "non-transfert base→instruct observé sur un cas",
            "poutine : argument corpus sans recensement complet",
            "juge cohérent en répétition, pas validé humainement",
        ],
        accent=ORANGE,
        fill=LIGHT_ORANGE,
        body_size=11.0,
    )
    add_card(
        slide,
        6.56,
        1.40,
        2.72,
        3.25,
        "Non démontré",
        [
            "nécessité de 9056",
            "généralité inter-modèle",
            "chaîne Interlab A1→A11 complète",
            "validité humaine des jugements Lodestar",
        ],
        accent=GRAY,
        fill=WHITE,
        body_size=11.4,
    )
    notes(
        slide,
        "Synthétiser l’état scientifique avec précision et sans effet de conclusion finale.",
        "Le projet est plus mature parce que ses revendications sont maintenant classées par niveau de preuve.",
        [
            "Cette slide est utile avec le PI : elle rend explicite ce que je peux défendre aujourd’hui et ce que je ne dois pas encore défendre.",
            "Le résultat 9056 est fort comme sufficiency demonstration.",
            "Les limites ne sont pas un appendice; elles structurent directement les prochaines expériences.",
        ],
        "1 min 30 s",
        "Révéler colonne par colonne : établi, prudent, non démontré. Ne pas passer trop vite sur la dernière colonne.",
    )

    # Slide 40
    slide = add_slide(prs)
    add_title(slide, "Prochaines décisions de recherche", "Prioriser les expériences qui augmentent le plus la force causale et la généralité.", "ACTE IV")
    add_card(slide, 0.78, 1.35, 4.02, 0.82, "1 · Ablation de 9056", "transformer la sufficience en test de nécessité", accent=ORANGE, fill=LIGHT_ORANGE)
    add_card(slide, 5.18, 1.35, 4.02, 0.82, "2 · Bras Gemma Scope", "tester si le classement et le steering se généralisent hors Qwen", accent=BLUE, fill=LIGHT_BLUE)
    add_card(slide, 0.78, 2.42, 4.02, 0.82, "3 · Circuits de 9056", "relier la feature à des mécanismes internes plutôt qu’au seul comportement", accent=TEAL, fill=LIGHT_TEAL)
    add_card(slide, 5.18, 2.42, 4.02, 0.82, "4 · Grille layer × width", "compléter la carte de santé SAE si les checkpoints existent", accent=GREEN, fill=LIGHT_GREEN)
    add_text(slide, "Critère de priorité", 0.96, 3.88, 2.10, 0.32, size=15, color=DARK, bold=True)
    add_text(slide, "Choisir l’expérience qui ferme le garde-fou le plus important avec le coût expérimental le plus raisonnable.", 0.96, 4.27, 7.95, 0.48, size=16, color=DARK)
    add_arrow(slide, 3.15, 4.12, 4.10, 4.12, color=TEAL, width=2.5)
    add_tag(slide, "priorité actuelle : 9056 ablation", 4.25, 3.95, 2.45, 0.42, fill=LIGHT_ORANGE, line=ORANGE, color=ORANGE, size=11)
    notes(
        slide,
        "Terminer sur des décisions de recherche, pas sur une conclusion finale.",
        "La prochaine étape la plus logique est l’ablation de 9056, car elle ferme le garde-fou principal du résultat positif.",
        [
            "L’ablation est prioritaire parce qu’elle transforme la revendication : de « clamping produit l’effet » vers « cette feature est nécessaire à l’effet ».",
            "Gemma Scope répond à la généralité inter-modèle, mais demande plus de surface expérimentale.",
            "Les circuits deviennent naturels une fois la feature 9056 stabilisée comme objet d’étude.",
            "La grille layer×width est utile, mais moins urgente que les contrôles causaux et la généralité.",
        ],
        "1 min 30 s",
        "Révéler les quatre options, puis le critère de priorité et l’étiquette finale. Terminer en ouvrant la discussion avec le PI.",
    )

    prs.save(str(DECK))
    write_notes_markdown(start_slide, NOTE_PAYLOADS)


if __name__ == "__main__":
    build_deck()
