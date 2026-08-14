# -*- coding: utf-8 -*-
"""Create truncated variants of sae.pptx for PowerPoint-open bisection."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation


def delete_slide(prs: Presentation, index: int) -> None:
    """Delete slide by zero-based index using python-pptx internals."""
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    prs.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def keep_first(src: str, dst: str, count: int) -> None:
    shutil.copy2(src, dst)
    prs = Presentation(dst)
    while len(prs.slides) > count:
        delete_slide(prs, len(prs.slides) - 1)
    prs.save(dst)


if __name__ == "__main__":
    for count in [24, 25, 26, 28, 30, 32, 34, 36, 38, 40]:
        keep_first("sae.pptx", f"diagnose_keep_{count}.pptx", count)
        print(f"diagnose_keep_{count}.pptx")
